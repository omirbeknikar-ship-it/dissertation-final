"""
build_presentation.py
Generates the final defense presentation as a .pptx file.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from pathlib import Path

ROOT   = Path(__file__).parent.parent
FIGS   = ROOT / "Outputs" / "generated_figures"
OUT    = ROOT / "Paper" / "defense_presentation.pptx"

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x38, 0x64)   # dark navy  – titles, title-slide bg
BLUE   = RGBColor(0x2E, 0x75, 0xB6)   # mid blue   – accents, dividers
RED    = RGBColor(0xC0, 0x00, 0x00)   # red        – key numbers / warnings
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)   # light grey – slide backgrounds
MGRAY  = RGBColor(0xBF, 0xBF, 0xBF)   # mid grey   – table rules

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helper utilities ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, italic=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=16, bold=False, italic=False,
             color=BLACK, align=PP_ALIGN.LEFT, space_before=6):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = _Pt(space_before)
    run = p.add_run()
    run.text   = text
    run.font.size   = _Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def content_slide(prs, title_text, subtitle_text=""):
    """White slide with navy top bar."""
    slide = prs.slides.add_slide(BLANK)
    # background
    add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
    # navy top bar
    add_rect(slide, 0, 0, 13.33, 1.1, fill=NAVY)
    # blue accent line
    add_rect(slide, 0, 1.1, 13.33, 0.06, fill=BLUE)
    # title
    add_text(slide, title_text, 0.4, 0.15, 12.5, 0.85,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_text(slide, subtitle_text, 0.4, 0.82, 12.5, 0.35,
                 size=14, italic=True, color=RGBColor(0xBF, 0xD7, 0xED),
                 align=PP_ALIGN.LEFT)
    return slide


def bullet_box(slide, items, l, t, w, h,
               size=17, color=BLACK, indent_marker="•"):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        # indent level
        if isinstance(item, tuple):
            txt, lvl = item
        else:
            txt, lvl = item, 0
        p.level = lvl
        run = p.add_run()
        marker = "  –" if lvl > 0 else indent_marker
        run.text = f"{marker}  {txt}"
        run.font.size  = Pt(size - lvl * 1.5)
        run.font.color.rgb = color
    return txb


def stat_box(slide, value, label, l, t, w=2.5, h=1.4,
             val_color=RED, bg=LGRAY):
    add_rect(slide, l, t, w, h, fill=bg, line=MGRAY)
    add_text(slide, value, l+0.1, t+0.08, w-0.2, 0.72,
             size=34, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, l+0.1, t+0.78, w-0.2, 0.55,
             size=13, color=NAVY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(s1, 0, 5.6, 13.33, 1.9, fill=BLUE)
add_rect(s1, 0, 5.55, 13.33, 0.08, fill=WHITE)

add_text(s1,
         "BRI's Impact on Kazakhstan–China Trade Balance:",
         0.6, 1.1, 12.0, 0.9,
         size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s1,
         "The Role of Strategic Mineral Exports",
         0.6, 1.95, 12.0, 0.7,
         size=26, bold=False, color=RGBColor(0xBF, 0xD7, 0xED),
         align=PP_ALIGN.LEFT)
add_rect(s1, 0.6, 2.75, 4.5, 0.05, fill=WHITE)
add_text(s1, "Nikar Omirbek", 0.6, 2.9, 6, 0.5,
         size=18, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s1, "Final Dissertation Defense  |  2026",
         0.6, 3.4, 7, 0.4,
         size=14, color=RGBColor(0xBF, 0xD7, 0xED), align=PP_ALIGN.LEFT)

add_text(s1,
         "Evidence from an Annual Bilateral Time-Series Analysis, 2000–2023",
         0.6, 6.0, 12.0, 0.5,
         size=13, italic=True,
         color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
s2 = content_slide(prs, "Presentation Overview")
items = [
    "The Puzzle — what the data shows at face value",
    "Research Question & Theory",
    "Data & Analytical Pipeline",
    "Regression Analysis & Multicollinearity Resolution",
    "Main Finding: TWFE Difference-in-Differences",
    "Robustness & Limitations",
    "Conclusions & Policy Implications",
]
bullet_box(s2, items, 0.8, 1.3, 11.5, 5.8, size=19, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PUZZLE  (stat boxes)
# ══════════════════════════════════════════════════════════════════════════════
s3 = content_slide(prs, "The Puzzle: BRI Was Supposed to Help Kazakhstan")

add_text(s3,
         "After BRI launched (2013), Kazakhstan's trade position with China deteriorated — "
         "despite record mineral exports.",
         0.5, 1.25, 12.3, 0.6, size=16, color=NAVY)

# stat boxes row 1
stat_box(s3, "+94.3%", "Imports from China\n(post-BRI avg.)", 0.5, 2.05, val_color=RED)
stat_box(s3, "+24.3%", "Exports to China\n(post-BRI avg.)", 3.2, 2.05, val_color=BLUE)
stat_box(s3, "+29.2%", "Mineral exports\n(post-BRI avg.)", 5.9, 2.05, val_color=BLUE)
stat_box(s3, "−37.8%", "Trade surplus\n(post-BRI avg.)", 8.6, 2.05, val_color=RED)

# 2023 callout
add_rect(s3, 0.5, 3.65, 11.8, 1.1, fill=RGBColor(0xFF,0xEB,0xEB), line=RED)
add_text(s3,
         "2023  —  First-ever bilateral trade DEFICIT with China:  −$2.01 billion",
         0.7, 3.75, 11.4, 0.55,
         size=19, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s3,
         "Imports $16.77bn  vs  Exports $14.76bn  |  Record mineral exports in the same year",
         0.7, 4.25, 11.4, 0.4,
         size=14, color=NAVY, align=PP_ALIGN.CENTER)

add_text(s3,
         "Core question:  Did mineral wealth provide any real leverage — or did asymmetric "
         "interdependence dominate?",
         0.5, 4.95, 12.3, 0.55,
         size=15, italic=True, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — RESEARCH QUESTION & THEORY
# ══════════════════════════════════════════════════════════════════════════════
s4 = content_slide(prs, "Research Question & Theoretical Framework")

add_text(s4,
         "Research Question",
         0.5, 1.25, 12, 0.45, size=18, bold=True, color=NAVY)
add_rect(s4, 0.5, 1.72, 12.2, 0.9, fill=LGRAY, line=BLUE)
add_text(s4,
         "How did the post-BRI period change the relationship between Kazakhstan's "
         "strategic mineral exports and its bilateral trade balance with China?",
         0.7, 1.78, 11.8, 0.78,
         size=17, italic=True, color=NAVY, align=PP_ALIGN.LEFT)

add_text(s4, "Theoretical Framework: Asymmetric Interdependence",
         0.5, 2.8, 12, 0.45, size=18, bold=True, color=NAVY)

cols = [
    ("Kazakhstan", ["Needs China as primary buyer", "Limited alternative markets",
                    "Concentrated mineral export basket", "→  High dependence"]),
    ("China", ["Multiple mineral suppliers globally", "Diversified import sources",
               "200× larger economy", "→  Low dependence on KZ"]),
    ("Prediction", ["BRI deepens the relationship", "Import surge outpaces exports",
                    "Mineral revenues insufficient", "→  Balance deteriorates"]),
]
for i, (hdr, pts) in enumerate(cols):
    x = 0.5 + i * 4.25
    clr = RED if i == 0 else (BLUE if i == 1 else NAVY)
    add_rect(s4, x, 3.35, 3.9, 0.45, fill=clr)
    add_text(s4, hdr, x+0.1, 3.38, 3.7, 0.38,
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s4, x, 3.8, 3.9, 2.8, fill=LGRAY, line=MGRAY)
    y = 3.88
    for pt in pts:
        add_text(s4, f"  {pt}", x+0.1, y, 3.7, 0.52, size=13, color=BLACK)
        y += 0.52


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DATA & PIPELINE
# ══════════════════════════════════════════════════════════════════════════════
s5 = content_slide(prs, "Data & Analytical Pipeline")

add_text(s5, "Dataset", 0.5, 1.25, 5.5, 0.4, size=17, bold=True, color=NAVY)
bullet_box(s5, [
    "Annual bilateral panel: 2000–2023  (n = 24)",
    "IMF DOTS — bilateral exports & imports",
    "World Bank WDI — GDP, exchange rates",
    "UN Comtrade — HS-level mineral exports",
    "Kazakhstan Bureau of National Statistics — cross-validation",
    "6 partner countries: CHN, RUS, DEU, UZB, TUR, USA",
], 0.5, 1.72, 5.8, 4.6, size=15, color=BLACK)

add_text(s5, "Reproducible GitHub Pipeline", 7.0, 1.25, 5.8, 0.4,
         size=17, bold=True, color=NAVY)
bullet_box(s5, [
    "36 Python scripts — fully automated",
    "Data collection → cleaning → analysis → figures",
    "Every table and chart generated by code",
    "Version-controlled, publicly accessible",
    "All results reproducible from raw data",
], 7.0, 1.72, 5.9, 3.2, size=15, color=BLACK)

add_rect(s5, 7.0, 5.0, 5.9, 1.6, fill=LGRAY, line=BLUE)
add_text(s5, "Validation: BNS-scraped figures vs panel",
         7.1, 5.08, 5.7, 0.38, size=13, bold=True, color=NAVY)
add_text(s5, "2023 exports: 0.3% discrepancy  |  2023 imports: 0.1% discrepancy",
         7.1, 5.45, 5.7, 0.38, size=13, color=BLACK)
add_text(s5, "Both well within the <10% validation threshold  ✓",
         7.1, 5.82, 5.7, 0.38, size=13, color=RGBColor(0x37, 0x86, 0x10))

add_rect(s5, 0.5, 5.0, 5.8, 0.05, fill=BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TRADE BALANCE FIGURE
# ══════════════════════════════════════════════════════════════════════════════
s6 = content_slide(prs, "Trade Balance Over Time (2000–2023)")
fig_path = FIGS / "fig_2_trade_balance.png"
if fig_path.exists():
    s6.shapes.add_picture(str(fig_path), Inches(0.5), Inches(1.25),
                          Inches(12.3), Inches(5.6))
add_text(s6,
         "The 2023 deficit is the first in recorded Kazakhstan–China trade history. "
         "A decade-long deterioration, not a single-year shock.",
         0.5, 6.9, 12.3, 0.5, size=12, italic=True, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — EXPORTS vs IMPORTS FIGURE
# ══════════════════════════════════════════════════════════════════════════════
s7 = content_slide(prs, "Import-Side Deepening: The Primary Mechanism")
fig_path = FIGS / "fig_1_exports_imports.png"
if fig_path.exists():
    s7.shapes.add_picture(str(fig_path), Inches(0.5), Inches(1.25),
                          Inches(8.5), Inches(5.3))
bullet_box(s7, [
    "Imports nearly quadrupled post-BRI",
    "Export growth modest (+24%)",
    "Mineral exports +29% — insufficient",
    "Import-side deepening is the\nprimary mechanism of deterioration",
], 9.2, 1.5, 3.85, 4.5, size=16, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — EMPIRICAL STRATEGY
# ══════════════════════════════════════════════════════════════════════════════
s8 = content_slide(prs, "Empirical Strategy")

rows = [
    ("OLS Regression\n(HAC s.e.)", "Tests whether post-BRI × minerals interaction\nis negative", "Primary coefficient of interest"),
    ("Multicollinearity\nResolution", "GDP levels → GDP growth rates\nVIF: 236 → 10.4", "Scheme B = primary specification"),
    ("288-Spec Grid", "Vary mineral measure, lags, estimator, BRI proxy", "Robustness to specification choice"),
    ("Leave-One-Out", "Remove each year; test coefficient stability", "Identifies influential observations"),
    ("Influence Diagnostics", "Cook's distance, leverage, studentised residuals", "Flags 2023 as extreme outlier"),
    ("TWFE DiD", "Kazakhstan vs 5 control partners\nPartner + year fixed effects", "Primary causal identification"),
    ("Synthetic Control", "Pre-period counterfactual trajectory", "Complementary causal check"),
]

y = 1.3
add_rect(s8, 0.4, y, 2.8, 0.38, fill=NAVY)
add_rect(s8, 3.3, y, 5.2, 0.38, fill=NAVY)
add_rect(s8, 8.6, y, 4.3, 0.38, fill=NAVY)
for hdr, x, w in [("Method", 0.4, 2.8), ("What it does", 3.3, 5.2), ("Purpose", 8.6, 4.3)]:
    add_text(s8, hdr, x+0.1, y+0.04, w-0.2, 0.3,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
y += 0.38
for i, (method, does, purpose) in enumerate(rows):
    bg = LGRAY if i % 2 == 0 else WHITE
    h = 0.72
    add_rect(s8, 0.4, y, 2.8, h, fill=bg, line=MGRAY)
    add_rect(s8, 3.3, y, 5.2, h, fill=bg, line=MGRAY)
    add_rect(s8, 8.6, y, 4.3, h, fill=bg, line=MGRAY)
    add_text(s8, method, 0.5, y+0.05, 2.6, h-0.1, size=12, bold=True, color=NAVY)
    add_text(s8, does,   3.4, y+0.05, 5.0, h-0.1, size=12, color=BLACK)
    add_text(s8, purpose,8.7, y+0.05, 4.1, h-0.1, size=12, italic=True, color=BLUE)
    y += h


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MULTICOLLINEARITY
# ══════════════════════════════════════════════════════════════════════════════
s9 = content_slide(prs, "Multicollinearity: Problem & Resolution")

add_text(s9,
         "Original model used log(KZ GDP) + log(CN GDP) as separate regressors. "
         "Both trend upward nearly every year → near-perfect collinearity.",
         0.5, 1.25, 12.3, 0.6, size=16, color=NAVY)

schemes = [
    ("Scheme A\n(Pre-revision)", "log(KZ GDP)\nlog(CN GDP)", "236.3", "❌  Severe", RED),
    ("Scheme B\n(Primary  ✓)", "d.log(KZ GDP)\nd.log(CN GDP)", "10.4", "✓  Meets threshold", RGBColor(0x37,0x86,0x10)),
    ("Scheme C\n(Robustness)", "log(CN/KZ GDP)\nratio", "33.4", "⚠  Improved; not ideal", RGBColor(0xFF,0x7F,0x00)),
]
for i, (name, vars_, vif, verdict, clr) in enumerate(schemes):
    x = 0.5 + i * 4.2
    add_rect(s9, x, 1.95, 3.9, 0.5, fill=clr)
    add_text(s9, name, x+0.1, 1.97, 3.7, 0.46,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s9, x, 2.45, 3.9, 3.3, fill=LGRAY, line=MGRAY)
    add_text(s9, "Variables:", x+0.15, 2.5, 3.6, 0.3, size=12, bold=True, color=NAVY)
    add_text(s9, vars_, x+0.15, 2.8, 3.6, 0.7, size=13, italic=True, color=BLACK)
    add_text(s9, "Max VIF:", x+0.15, 3.55, 1.5, 0.3, size=12, bold=True, color=NAVY)
    add_text(s9, vif, x+1.7, 3.52, 2.0, 0.38, size=24, bold=True, color=clr, align=PP_ALIGN.RIGHT)
    add_text(s9, verdict, x+0.15, 4.05, 3.6, 0.55, size=13, bold=True, color=clr)

add_rect(s9, 0.5, 5.85, 12.2, 0.6, fill=RGBColor(0xE2,0xEF,0xDA), line=RGBColor(0x37,0x86,0x10))
add_text(s9,
         "Scheme B is the primary specification: GDP growth rates are stationary, "
         "theoretically justified, and achieve VIF = 10.4.",
         0.65, 5.92, 12.0, 0.46,
         size=14, bold=True, color=RGBColor(0x37,0x5C,0x00))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — REGRESSION RESULTS
# ══════════════════════════════════════════════════════════════════════════════
s10 = content_slide(prs, "Regression Results: Post-BRI × Minerals Interaction")

add_text(s10, "Key coefficient:  Post-BRI × Mineral Exports  (β₃)",
         0.5, 1.25, 12, 0.42, size=17, bold=True, color=NAVY)

tbl_data = [
    ("Specification",           "N",  "Interaction β",  "HAC SE",  "p-value",  "Max VIF"),
    ("A: GDP levels (pre-rev)", "24", "−3.220",         "1.043",   "0.002",    "236.3"),
    ("B: GDP growth ✓ primary","23", "−2.489",         "0.974",   "0.011",    "10.4"),
    ("B: excl. 2022–2023",      "21", "−0.012",         "0.283",   "0.966",    "10.4"),
    ("B: excl. 2023 only",      "22", "+0.065",         "0.239",   "0.785",    "10.4"),
    ("C: Gravity ratio",        "24", "−2.419",         "1.131",   "0.033",    "33.4"),
    ("C: excl. 2022–2023",      "22", "+0.102",         "0.432",   "0.813",    "33.4"),
]

col_w = [4.2, 0.7, 1.5, 1.2, 1.2, 1.2]
col_x = [0.5]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

y = 1.75
for ri, row in enumerate(tbl_data):
    if ri == 0:
        bg = NAVY; fc = WHITE; fsize = 12; fb = True
    elif ri == 2:
        bg = RGBColor(0xE2,0xEF,0xDA); fc = BLACK; fsize = 13; fb = True
    else:
        bg = LGRAY if ri % 2 else WHITE; fc = BLACK; fsize = 12; fb = False
    rh = 0.46 if ri > 0 else 0.38
    for ci, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        add_rect(s10, cx, y, cw, rh, fill=bg, line=MGRAY)
        clr = (RED if (ri > 0 and ci in [2,4] and "−" in str(cell) and ri not in [3,4,6,7]) else fc)
        add_text(s10, cell, cx+0.05, y+0.04, cw-0.1, rh-0.08,
                 size=fsize, bold=fb, color=clr,
                 align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
    y += rh

add_text(s10,
         "⚠  The interaction disappears when 2022–2023 are excluded → result is fragile. "
         "Primary evidence comes from the DiD (next slide).",
         0.5, 6.7, 12.3, 0.5, size=13, italic=True, color=RED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — INFLUENCE DIAGNOSTICS
# ══════════════════════════════════════════════════════════════════════════════
s11 = content_slide(prs, "Influence Diagnostics: The Role of 2023")
fig_path = FIGS / "fig_5_diagnostics.png"
if fig_path.exists():
    s11.shapes.add_picture(str(fig_path), Inches(0.5), Inches(1.25),
                           Inches(8.2), Inches(5.0))

add_rect(s11, 9.0, 1.4, 3.9, 1.3, fill=RGBColor(0xFF,0xEB,0xEB), line=RED)
add_text(s11, "2023", 9.1, 1.45, 3.7, 0.42, size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s11, "Cook's D = 3.90\n23× the threshold (4/n = 0.167)",
         9.1, 1.85, 3.7, 0.75, size=13, color=BLACK, align=PP_ALIGN.CENTER)

bullet_box(s11, [
    "Excluding any year except 2023\n→ interaction stays negative & significant",
    "Excluding 2023 alone\n→ interaction loses significance",
    "R² jumps from 0.765 → 0.960\nwhen 2023 removed",
    "This is reported transparently\nin the paper",
], 9.0, 2.85, 3.9, 3.8, size=14, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — DiD MAIN FINDING
# ══════════════════════════════════════════════════════════════════════════════
s12 = content_slide(prs, "Main Finding: TWFE Difference-in-Differences",
                    "Is the deterioration China-specific, or a global trend?")

add_text(s12,
         "Design:  Stack Kazakhstan + 5 control partners (RUS, DEU, UZB, TUR, USA)  "
         "→  144 obs.  |  Partner FE + Year FE  |  HAC s.e.",
         0.5, 1.25, 12.3, 0.5, size=15, color=NAVY)

# big result box
add_rect(s12, 0.5, 1.85, 5.8, 2.3, fill=NAVY)
add_text(s12, "DiD Coefficient", 0.7, 1.92, 5.4, 0.45,
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s12, "−0.305", 0.7, 2.32, 5.4, 0.95,
         size=58, bold=True, color=RGBColor(0xFF,0xCC,0x00), align=PP_ALIGN.CENTER)
add_text(s12, "p = 0.0002   |   HAC SE = 0.083", 0.7, 3.25, 5.4, 0.42,
         size=15, color=RGBColor(0xBF,0xD7,0xED), align=PP_ALIGN.CENTER)

add_rect(s12, 6.6, 1.85, 6.2, 2.3, fill=LGRAY, line=BLUE)
add_text(s12, "What this means:", 6.75, 1.92, 5.9, 0.38,
         size=14, bold=True, color=NAVY)
bullet_box(s12, [
    "Kazakhstan's balance ratio fell ~30.5 pp more\nthan control partners post-2013",
    "After controlling for all common time trends\nand partner-level differences",
    "Not driven by oil prices, global macro,\nor general China trade patterns",
    "→  The deterioration is CHINA-SPECIFIC",
], 6.6, 2.3, 6.2, 1.85, size=14, color=BLACK)

# robustness check row
add_rect(s12, 0.5, 4.25, 5.8, 1.0, fill=LGRAY, line=MGRAY)
add_text(s12, "China vs Russia only:", 0.65, 4.32, 5.5, 0.35,
         size=13, bold=True, color=NAVY)
add_text(s12, "DiD = −0.151  (p = 0.008)  —  robust to most similar partner comparison",
         0.65, 4.62, 5.5, 0.55, size=13, color=BLACK)

add_rect(s12, 6.6, 4.25, 6.2, 1.0, fill=RGBColor(0xE2,0xEF,0xDA), line=RGBColor(0x37,0x86,0x10))
add_text(s12, "✓  This is the strongest identification result in the paper.",
         6.75, 4.55, 5.9, 0.55, size=14, bold=True, color=RGBColor(0x37,0x5C,0x00))

# DiD figure
fig_path = FIGS / "fig_did_event_study.png"
if fig_path.exists():
    s12.shapes.add_picture(str(fig_path), Inches(0.5), Inches(5.4),
                           Inches(12.3), Inches(1.9))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — SYNTHETIC CONTROL
# ══════════════════════════════════════════════════════════════════════════════
s13 = content_slide(prs, "Synthetic Counterfactual: What If BRI Hadn't Happened?")
fig_path = FIGS / "fig_synthetic_control.png"
if fig_path.exists():
    s13.shapes.add_picture(str(fig_path), Inches(0.5), Inches(1.25),
                           Inches(8.5), Inches(5.3))
bullet_box(s13, [
    "Pre-period (2000–2012) dynamics used\nto project a synthetic counterfactual",
    "Gap negative in 8 of 11 post-BRI years",
    "MSPE ratio = 9.53\n(post-period 9× harder to fit)",
    "Consistent direction with DiD",
    "Not statistically significant alone\n→ treated as complementary evidence",
], 9.2, 1.4, 3.85, 5.0, size=15, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — TRIPLE CONCORDANCE
# ══════════════════════════════════════════════════════════════════════════════
s14 = content_slide(prs, "Evidence Summary: Triple-Concordance Assessment")

add_text(s14,
         "All methods point in the same direction (deterioration). Strength varies.",
         0.5, 1.25, 12.3, 0.42, size=16, color=NAVY)

methods = [
    ("Regression (Scheme B primary)",    "Negative (−2.49)", "p = 0.011", "Driven by 2022–23;  fragile",                        "⚠", RGBColor(0xFF,0x7F,0x00)),
    ("Regression (Scheme C robustness)", "Negative (−2.42)", "p = 0.033", "Same fragility pattern",                             "⚠", RGBColor(0xFF,0x7F,0x00)),
    ("TWFE DiD — China vs all partners", "Negative (−0.305)","p = 0.0002","China-specific;  strongest ID",                      "✓✓",RGBColor(0x37,0x86,0x10)),
    ("TWFE DiD — China vs Russia",       "Negative (−0.151)","p = 0.008", "Robust to similar partner",                          "✓", RGBColor(0x37,0x86,0x10)),
    ("Synthetic control gap",            "Negative (−0.014)","p = 0.857", "Consistent direction; not significant",              "→", BLUE),
    ("Sanctions robustness excl 22–23",  "Insignificant",    "p = 0.966", "Alternative channel cannot be ruled out",            "⚠", RGBColor(0xFF,0x7F,0x00)),
]

y = 1.75
col_x2 = [0.4, 3.3, 5.5, 7.0, 11.9, 12.5]
col_w2 = [2.9, 2.2, 1.5, 4.9, 0.6]
hdrs = ["Method", "Direction", "p-value", "Notes", ""]
add_rect(s14, 0.4, y, 12.5, 0.4, fill=NAVY)
for hdr, cx, cw in zip(hdrs, col_x2, col_w2):
    add_text(s14, hdr, cx+0.05, y+0.04, cw-0.1, 0.32,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
y += 0.4
for i, (m, d, p, n, sym, clr) in enumerate(methods):
    bg = LGRAY if i % 2 == 0 else WHITE
    rh = 0.52
    for cx, cw in zip(col_x2[:-1], col_w2[:-1]):
        add_rect(s14, cx, y, cw, rh, fill=bg, line=MGRAY)
    vals = [m, d, p, n]
    for ci, (val, cx, cw) in enumerate(zip(vals, col_x2, col_w2)):
        add_text(s14, val, cx+0.05, y+0.06, cw-0.1, rh-0.1,
                 size=12, color=(clr if ci == 1 else BLACK))
    add_text(s14, sym, 12.1, y+0.06, 0.7, rh-0.1,
             size=16, bold=True, color=clr, align=PP_ALIGN.CENTER)
    y += rh

add_rect(s14, 0.4, y+0.1, 12.5, 0.55, fill=RGBColor(0xE2,0xEF,0xDA), line=RGBColor(0x37,0x86,0x10))
add_text(s14,
         "Conclusion:  Evidence is strongly consistent with asymmetric interdependence deepening post-BRI.  "
         "Causal attribution to BRI specifically requires caution.",
         0.55, y+0.17, 12.2, 0.38, size=13, bold=True, color=RGBColor(0x37,0x5C,0x00))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — LIMITATIONS
# ══════════════════════════════════════════════════════════════════════════════
s15 = content_slide(prs, "Limitations: What This Study Cannot Establish")

cannot = [
    "BRI caused the trade balance deterioration",
    "The regression interaction is a stable structural parameter",
    "The 2022–2023 import surge is purely BRI-driven (sanctions parallel-import channel exists)",
    "Mineral-specific effects can be fully isolated from energy-export dynamics (HS-27 data unavailable)",
]
can = [
    "The post-BRI trade balance deterioration is China-specific (TWFE DiD, p = 0.0002)",
    "Strategic mineral export growth did not offset the import-side surge",
    "The pattern is strongly consistent with asymmetric interdependence theory",
    "The methodology is transparent, reproducible, and diagnostically honest",
]

add_rect(s15, 0.4, 1.25, 6.1, 0.42, fill=RED)
add_text(s15, "Cannot claim", 0.5, 1.29, 5.9, 0.34,
         size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s15, 0.4, 1.67, 6.1, 4.3, fill=RGBColor(0xFF,0xEB,0xEB), line=RED)
bullet_box(s15, cannot, 0.55, 1.72, 5.85, 4.2, size=15, color=BLACK)

add_rect(s15, 6.8, 1.25, 6.1, 0.42, fill=RGBColor(0x37,0x86,0x10))
add_text(s15, "Can claim", 6.9, 1.29, 5.9, 0.34,
         size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s15, 6.8, 1.67, 6.1, 4.3, fill=RGBColor(0xE2,0xEF,0xDA), line=RGBColor(0x37,0x86,0x10))
bullet_box(s15, can, 6.95, 1.72, 5.85, 4.2, size=15, color=BLACK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — POLICY IMPLICATIONS
# ══════════════════════════════════════════════════════════════════════════════
s16 = content_slide(prs, "Policy Implications")

add_text(s16,
         "Kazakhstan is not a unique case — findings are relevant to all resource-rich BRI partners.",
         0.5, 1.25, 12.3, 0.48, size=16, italic=True, color=NAVY)

policies = [
    ("Monitor trade balance composition",
     "Aggregate trade volume growth can mask a deteriorating bilateral position."),
    ("Track import-side dependence",
     "94% import growth signals deepening structural reliance on Chinese goods."),
    ("Prioritise value-added export upgrading",
     "Raw mineral exports are price-sensitive and insufficient to offset import surges."),
    ("Evaluate 2023 as a structural inflection",
     "First-ever deficit may mark a new regime — not a one-time statistical anomaly."),
    ("Diversify export destinations",
     "Concentration risk in a single large buyer amplifies asymmetric vulnerability."),
]

y = 1.82
for i, (title, desc) in enumerate(policies):
    bg = LGRAY if i % 2 == 0 else WHITE
    add_rect(s16, 0.4, y, 0.55, 0.8, fill=BLUE)
    add_text(s16, str(i+1), 0.4, y+0.18, 0.55, 0.44,
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s16, 0.95, y, 11.95, 0.8, fill=bg, line=MGRAY)
    add_text(s16, title, 1.1, y+0.04, 11.6, 0.36, size=14, bold=True, color=NAVY)
    add_text(s16, desc,  1.1, y+0.38, 11.6, 0.36, size=13, color=BLACK)
    y += 0.8


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — CONCLUSIONS
# ══════════════════════════════════════════════════════════════════════════════
s17 = content_slide(prs, "Conclusions")

concls = [
    ("Descriptive finding",
     "Imports +94% vs exports +24%. Trade surplus fell 37.8%. First deficit in 2023."),
    ("Regression finding",
     "Negative post-BRI × minerals interaction (β = −2.49, p = 0.011) — but fragile to exclusion of 2022–2023."),
    ("Strongest finding",
     "TWFE DiD: Kazakhstan's balance deteriorated 30.5 pp more than control partners (p = 0.0002). "
     "Deterioration is China-specific."),
    ("Interpretation",
     "Strongly consistent with asymmetric interdependence. "
     "Mineral wealth insufficient to offset China's import-capacity advantage."),
    ("Broader relevance",
     "Applicable to all resource-rich BRI partner economies facing structural asymmetry "
     "with China as their dominant trading partner."),
]

y = 1.35
for hdr, body in concls:
    add_rect(s17, 0.4, y, 2.5, 0.85, fill=NAVY)
    add_text(s17, hdr, 0.5, y+0.15, 2.35, 0.56,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s17, 2.9, y, 10.0, 0.85, fill=LGRAY, line=MGRAY)
    add_text(s17, body, 3.05, y+0.1, 9.75, 0.65, size=14, color=BLACK)
    y += 0.9

add_rect(s17, 0.4, y+0.1, 12.5, 0.55, fill=BLUE)
add_text(s17,
         "Transparent methodology  •  Reproducible pipeline  •  Honest reporting of fragility",
         0.5, y+0.2, 12.3, 0.35,
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
s18 = prs.slides.add_slide(BLANK)
add_rect(s18, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(s18, 0, 5.5, 13.33, 2.0, fill=BLUE)
add_rect(s18, 0, 5.45, 13.33, 0.08, fill=WHITE)

add_text(s18, "Thank You", 0.6, 1.4, 12.0, 1.0,
         size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s18, "Questions & Discussion",
         0.6, 2.5, 12.0, 0.6,
         size=26, color=RGBColor(0xBF,0xD7,0xED), align=PP_ALIGN.CENTER)
add_rect(s18, 3.5, 3.25, 6.3, 0.06, fill=WHITE)

add_text(s18,
         "Key number to remember:   DiD = −0.305   (p = 0.0002)   — China-specific",
         0.6, 3.5, 12.0, 0.55,
         size=16, italic=True, color=RGBColor(0xFF,0xCC,0x00), align=PP_ALIGN.CENTER)

add_text(s18, "Nikar Omirbek  |  2026",
         0.6, 5.8, 12.0, 0.45,
         size=15, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s18,
         "Full pipeline: github.com/omirbeknikar-ship-it/dissertation-final",
         0.6, 6.25, 12.0, 0.4,
         size=13, color=RGBColor(0xBF,0xD7,0xED), align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"Saved: {OUT}  ({OUT.stat().st_size:,} bytes)  —  {len(prs.slides)} slides")
